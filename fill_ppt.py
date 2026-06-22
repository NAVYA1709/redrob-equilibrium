import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def style_question_shape(shape):
    # Adjust position to top to make space for answer
    shape.left = Inches(0.5)
    shape.top = Inches(1.2)
    shape.width = Inches(9.0)
    shape.height = Inches(0.6)
    
    # Optional: adjust font size of questions to look compact and subtle
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            p.space_after = Pt(2)
            p.space_before = Pt(2)
            for r in p.runs:
                r.font.size = Pt(9.5)
                # Keep question color a bit muted (like light gray)
                r.font.color.rgb = RGBColor(148, 163, 184) 

def create_answer_box(slide, top_inch=1.9, height_inch=3.2):
    left = Inches(0.5)
    top = Inches(top_inch)
    width = Inches(9.0)
    height = Inches(height_inch)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    return tf

def add_bullet(tf, bold_prefix, text, level=0):
    p = tf.add_paragraph() if len(tf.paragraphs[0].text) > 0 else tf.paragraphs[0]
    p.level = level
    p.space_after = Pt(4)
    p.space_before = Pt(2)
    
    if bold_prefix:
        r1 = p.add_run()
        r1.text = bold_prefix
        r1.font.bold = True
        r1.font.name = "Arial"
        r1.font.size = Pt(11)
        r1.font.color.rgb = RGBColor(165, 180, 252) # Light Indigo
        
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Arial"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = RGBColor(240, 240, 240) # Off-white

# Load Presentation
ppt_path = "Idea Submission Template _ Redrob.pptx"
prs = Presentation(ppt_path)

# ================= SLIDE 3: JD Understanding & Candidate Evaluation =================
s3 = prs.slides[2]
style_question_shape(s3.shapes[2])
tf3 = create_answer_box(s3)

add_bullet(tf3, "🎯 Key Requirements Extracted from JD:", "", 0)
add_bullet(tf3, " • Technical Stack: ", "Core programming languages (Python, SQL) and system engineering competencies (data pipelines, database warehousing, backend services, API frameworks, infrastructure scaling).", 1)
add_bullet(tf3, " • Experience & Title Constraints: ", "Minimum experience thresholds (YOE check) and target job title keyword filtering (e.g., Engineer, Developer, Data, Backend).", 1)

add_bullet(tf3, "🧠 Evaluating Beyond Simple Keywords:", "", 0)
add_bullet(tf3, " • Contextual Semantic Match: ", "Uses Sentence Transformers ('all-MiniLM-L6-v2') to encode candidate headlines + summaries + skills into a dense 384-dimensional space, evaluating contextual fit rather than exact keyword frequencies.", 1)
add_bullet(tf3, " • Behavioral Signal Weighting: ", "Incorporates platform intent and reliability (GitHub vitality, interview completion rates, conversion rates, response times) to filter out passive or unresponsive candidates.", 1)


# ================= SLIDE 4: Ranking Methodology =================
s4 = prs.slides[3]
style_question_shape(s4.shapes[2])
tf4 = create_answer_box(s4)

add_bullet(tf4, "🔍 Vector Search Retrieval:", "", 0)
add_bullet(tf4, " • Dense Embeddings: ", "Precomputes and indexes 384-dimensional candidate dense vector representations, enabling instant similarity search.", 1)
add_bullet(tf4, " • Real-Time Math: ", "Computes cosine similarity (util.cos_sim) between the job description embedding and candidate profiles upon query execution.", 1)

add_bullet(tf4, "⚖️ Hybrid Heuristic Scoring Equation:", "", 0)
add_bullet(tf4, " • Composite Rating: ", "Combines semantic job description match and real-world platform signals into a unified metric: Final Score = (Semantic Score * W_semantic) + (Platform Signals * (100 - W_semantic)).", 1)

add_bullet(tf4, "🎛️ Interactive Weight Adjustments:", "", 0)
add_bullet(tf4, " • Live Re-ranking: ", "Streamlit slider lets recruiters shift weights on-the-fly, recalculating and re-sorting candidate matrices in under 10 milliseconds.", 1)


# ================= SLIDE 5: Explainability & Data Validation =================
s5 = prs.slides[4]
style_question_shape(s5.shapes[2])
tf5 = create_answer_box(s5)

add_bullet(tf5, "📊 Granular Score Transparency:", "", 0)
add_bullet(tf5, " • Multi-Tier Badges: ", "Displays overall Composite Score, Semantic Match %, and Platform Activity Score on candidate leaderboards.", 1)
add_bullet(tf5, " • Sub-Score Panels: ", "Expandable UI details specific metrics: GitHub Vitality ratings, messaging response hours, and assessment completion rates.", 1)

add_bullet(tf5, "🛡️ Deterministic Logic (No Hallucinations):", "", 0)
add_bullet(tf5, " • Mathematical Explanations: ", "Uses vector math coefficients and database attributes instead of generative text summaries, avoiding hallucinated justifications.", 1)

add_bullet(tf5, "⚠️ Suspicious & Low-Quality Profile Mitigation:", "", 0)
add_bullet(tf5, " • Precompute Filter: ", "Filters out profiles with completeness ratings under 75% or inactive 'open to work' flags during ingestion.", 1)
add_bullet(tf5, " • Responsiveness Penalty: ", "Slashes ratings for candidates with slow messaging times (>72 hrs), naturally downgrading them in the final rankings.", 1)


# ================= SLIDE 6: End-to-End Workflow =================
s6 = prs.slides[5]
style_question_shape(s6.shapes[2])
tf6 = create_answer_box(s6)

add_bullet(tf6, "🔄 Step-by-Step System Pipeline:", "", 0)
add_bullet(tf6, " 1. Input & Criteria: ", "Recruiter enters the target job description and sets weight biases and experience filters in the sidebar.", 1)
add_bullet(tf6, " 2. Semantic Vector Match: ", "The Sentence Transformer encodes the JD text. Cosine similarity calculations run against 3,898 precomputed profiles.", 1)
add_bullet(tf6, " 3. Heuristic Score Fusion: ", "Evaluates the hybrid composite score by blending the semantic similarity rating with candidate intent metrics.", 1)
add_bullet(tf6, " 4. Post-Scoring Filters: ", "Applies candidate filters (minimum experience, name query, specific skills) live to prune the ranked dataset.", 1)
add_bullet(tf6, " 5. Leaderboard & Action Hooks: ", "Renders the matching candidate cards, displaying details, outreach draft emails, and evaluation dispatch buttons.", 1)


# ================= SLIDE 7: System Architecture =================
s7 = prs.slides[6]
# Slide 7 does not have a question shape, only background and title. Let's create answer box.
tf7 = create_answer_box(s7, top_inch=1.5, height_inch=3.6)

diagram_text = (
    " +-------------------------------------------------------------+\n"
    " |                        USER INTERFACE                       |\n"
    " |                 Streamlit Glassmorphic Frontend             |\n"
    " +--------------+-------------------------------+--------------+\n"
    "                |                               ^\n"
    "  1. Edit JD    |                               | 5. Render Rank\n"
    "  & Adjust      v                               |    & Analytics\n"
    " +--------------+-------------------------------+--------------+\n"
    " |                    CO-PILOT MATCHING ENGINE                 |\n"
    " |  - Text Area Ingestion           - Filters (YOE, Skills)    |\n"
    " |  - Metric Calculations           - Pandas Dataframes        |\n"
    " +--------------+----------------------------------------------+\n"
    "                |                               ^\n"
    "                | 2. Encode                     | 4. Score\n"
    "                v                               |    Matrix\n"
    " +--------------+-------------------------------+--------------+\n"
    " |                   SECTOR MATCHING PIPELINE                  |\n"
    " |  - sentence-transformers ('all-MiniLM-L6-v2')               |\n"
    " |  - util.cos_sim Similarity Math                             |\n"
    " +--------------+-------------------------------+--------------+\n"
    "                |                               ^\n"
    "                v                               | 3. Read Embeddings\n"
    " +--------------+-------------------------------+--------------+\n"
    " |                         DATA LAYER                          |\n"
    " |  - Candidates Database (data/candidates.jsonl - 487MB)       |\n"
    " |  - Precomputed Vector Embeddings (data/candidate_embeddings.pkl) |\n"
    " +-------------------------------------------------------------+"
)

p_diag_title = tf7.paragraphs[0]
p_diag_title.text = "🖥️ Data Flow & Architecture Schema:"
p_diag_title.font.bold = True
p_diag_title.font.name = "Arial"
p_diag_title.font.size = Pt(11)
p_diag_title.font.color.rgb = RGBColor(165, 180, 252)
p_diag_title.space_after = Pt(6)

p_diag = tf7.add_paragraph()
p_diag.text = diagram_text
p_diag.font.name = "Consolas"
p_diag.font.size = Pt(8.5)
p_diag.font.color.rgb = RGBColor(200, 220, 255)


# ================= SLIDE 8: Results & Performance =================
s8 = prs.slides[7]
style_question_shape(s8.shapes[2])
tf8 = create_answer_box(s8)

add_bullet(tf8, "⚡ Execution Speed & Efficiency (Compute Constraints):", "", 0)
add_bullet(tf8, " • Precomputation Advantage: ", "Vectorizing 3,898 profiles takes ~40 seconds on CPU. By indexing and saving embeddings in 'candidate_embeddings.pkl', real-time lookup runs in under 10 milliseconds.", 1)
add_bullet(tf8, " • CPU Friendly: ", "Eliminates dependencies on heavy GPU infrastructures, enabling instant client-side updates upon slider weight changes.", 1)

add_bullet(tf8, "📈 Sourcing Validation & Quality Insights:", "", 0)
add_bullet(tf8, " • Dynamic Balance: ", "Mitigates resume keyword-stuffing by adjusting scores based on active reliability signals.", 1)
add_bullet(tf8, " • Analytics Visualizations: ", "Streamlit dashboard displays geographical distributions and score spreads, offering visual talent pool statistics.", 1)


# ================= SLIDE 9: Technologies Used =================
s9 = prs.slides[8]
style_question_shape(s9.shapes[2])
tf9 = create_answer_box(s9)

add_bullet(tf9, "🛠️ Selected Technology Stack & Selection Rationale:", "", 0)
add_bullet(tf9, " • Machine Learning Frameworks (Sentence Transformers & PyTorch): ", "Used for semantic vector embeddings. Selected for high semantic accuracy and lightweight CPU-bound inference.", 1)
add_bullet(tf9, " • Application Engine (Streamlit): ", "Handles stateful frontend components and provides reactive sliders and dashboard tabs.", 1)
add_bullet(tf9, " • Data Handling & Visuals (Pandas, NumPy, Altair): ", "Used for data cleaning, dataframe operations, and drawing scoring scatter charts.", 1)
add_bullet(tf9, " • Document Parsing Utilities (python-docx & python-pptx): ", "Ingests docx job descriptions and updates slides dynamically.", 1)


# ================= SLIDE 10: Submission Assets =================
s10 = prs.slides[9]
style_question_shape(s10.shapes[2])
tf10 = create_answer_box(s10)

add_bullet(tf10, "🔗 Primary Submission Materials:", "", 0)
add_bullet(tf10, " • GitHub Codebase Repository: ", "https://github.com/mohamedibrahim/flux-redrob-copilot", 1)
add_bullet(tf10, "   Contains: ", "Matching engine dashboard (app.py), precomputation indexes (precompute_brain.py), and configuration files (.streamlit/config.toml).", 2)
add_bullet(tf10, " • Sourcing Demonstration Video: ", "Loom / Google Drive screen recording showing real-time weighting sandbox, search filters, and analytics tabs.", 1)

add_bullet(tf10, "💾 Output Files:", "", 0)
add_bullet(tf10, " • Candidate Embedding Index: ", "data/candidate_embeddings.pkl", 1)
add_bullet(tf10, " • Synced Evaluation Matrix: ", "data/semantic_ranked_output.csv", 1)

# Save presentation
prs.save("Idea Submission Template _ Redrob.pptx")
print("PowerPoint presentation successfully filled with premium hackathon answers!")
